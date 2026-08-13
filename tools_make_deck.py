"""Build the client briefing deck for CISCLD Ascend.

  py -3.12 tools_make_deck.py

Writes data/CISCLD-Ascend-client-briefing.pptx.

Every figure quoted in the deck is an output of this repository's own engines
against the 547-VM reference estate, or a rate read from a vendor's live pricing
API. Nothing here is a marketing number, which is the point the deck itself
makes: a tool whose selling line is "we will show you the numbers that argue
against the migration" cannot have an invented number in its own briefing.

Regenerate after changing anything it quotes. The figures live in FACTS below so
there is exactly one place to correct.
"""

import pathlib
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = pathlib.Path("data") / "CISCLD-Ascend-client-briefing.pptx"
URL = "onpremtocloud.streamlit.app"

# --------------------------------------------------------------------------
# Palette, matching core/brand.py so the deck and the product are one thing
# --------------------------------------------------------------------------
VOID = RGBColor(0x08, 0x0C, 0x17)
PLATE = RGBColor(0x0E, 0x14, 0x26)
ACCENT = RGBColor(0x3B, 0x6F, 0xD4)
ACCENT_2 = RGBColor(0x7A, 0x5B, 0xC4)
POSITIVE = RGBColor(0x2F, 0x9E, 0x68)
WARNING = RGBColor(0xE0, 0x8A, 0x2E)
NEGATIVE = RGBColor(0xC4, 0x48, 0x5E)
INK = RGBColor(0x1E, 0x24, 0x2E)
INK_DIM = RGBColor(0x5A, 0x63, 0x72)
PAPER = RGBColor(0xE8, 0xED, 0xF7)
PAPER_DIM = RGBColor(0x9F, 0xAC, 0xC4)
RULE = RGBColor(0xDD, 0xE3, 0xEE)
CARD = RGBColor(0xF5, 0xF7, 0xFB)

SANS = "Segoe UI"
SANS_LIGHT = "Segoe UI Light"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)                     # page margin
CONTENT_W = W - 2 * M

# --------------------------------------------------------------------------
# Every quoted figure, in one place
# --------------------------------------------------------------------------
FACTS = {
    "vcpu_rate": "$0.0460",
    "migrate_coverage": "46%",
    "rightsizing": "16%",
    "mean_cpu": "15%",
    "avs_hosts": "31",
    "avs_saving": "5%",
    "pages": "16",
    "engines": "12",
}


# --------------------------------------------------------------------------
# Primitives
# --------------------------------------------------------------------------
def slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = VOID if dark else RGBColor(0xFF, 0xFF, 0xFF)
    return s


def rect(s, x, y, w, h, colour, shape=MSO_SHAPE.RECTANGLE):
    box = s.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = colour
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0, gap=0):
    """runs = [(string, size_pt, bold, colour, font)] -- one paragraph each.

    ``gap`` is the space in points after every paragraph but the last. Stacking
    related lines in one frame rather than positioning each in its own textbox is
    what keeps a two-line heading from landing on top of the body beneath it.
    """
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (body, size, bold, colour, font) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if gap and i < len(runs) - 1:
            p.space_after = Pt(gap)
        r = p.add_run()
        r.text = body
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = font
    return box


def eyebrow(s, label, colour=ACCENT):
    text(s, M, Inches(0.62), CONTENT_W, Inches(0.3),
         [(label.upper(), 11, True, colour, SANS)])


def heading(s, title, sub="", dark=False):
    text(s, M, Inches(1.02), CONTENT_W, Inches(0.85),
         [(title, 32, True, PAPER if dark else INK, SANS)])
    if sub:
        text(s, M, Inches(1.86), Inches(9.6), Inches(0.7),
             [(sub, 14, False, PAPER_DIM if dark else INK_DIM, SANS)], spacing=1.25)


def accent_bar(s, y=Inches(0.0), h=Inches(0.09)):
    rect(s, Inches(0), y, W * 0.42, h, ACCENT)
    rect(s, W * 0.42, y, W * 0.28, h, ACCENT_2)


def footer(s, n, dark=False):
    text(s, M, H - Inches(0.55), Inches(6), Inches(0.3),
         [(f"CISCLD Ascend  ·  {URL}", 9, False,
           PAPER_DIM if dark else INK_DIM, SANS)])
    text(s, W - M - Inches(1), H - Inches(0.55), Inches(1), Inches(0.3),
         [(str(n), 9, False, PAPER_DIM if dark else INK_DIM, SANS)],
         align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, body, accent=ACCENT, big=None):
    rect(s, x, y, w, h, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, w, Inches(0.055), accent)
    cy = y + Inches(0.34)
    if big:
        text(s, x + Inches(0.32), cy, w - Inches(0.5), Inches(0.62),
             [(big, 30, True, accent, MONO)])
        cy += Inches(0.68)
    # Title and body share one frame, so a heading that wraps pushes the body
    # down instead of sitting on top of it.
    text(s, x + Inches(0.32), cy, w - Inches(0.64), h - (cy - y) - Inches(0.3),
         [(title, 14, True, INK, SANS),
          (body, 11.5, False, INK_DIM, SANS)], spacing=1.22, gap=7)


def bullets(s, x, y, w, items, size=13.5, gap=Inches(0.62), colour=INK):
    for i, (head, body) in enumerate(items):
        yy = y + i * gap
        rect(s, x, yy + Inches(0.09), Inches(0.05), Inches(0.22), ACCENT)
        runs = [(head, size, True, colour, SANS)]
        if body:
            runs.append((body, 11.5, False, INK_DIM, SANS))
        text(s, x + Inches(0.24), yy, w - Inches(0.24), gap,
             runs, spacing=1.22, gap=5)


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------
def title_slide(prs):
    s = slide(prs, dark=True)
    accent_bar(s)
    rect(s, Inches(0), Inches(0), Inches(0.075), H, ACCENT)
    text(s, M, Inches(2.05), Inches(10), Inches(0.4),
         [("C I S C L D", 13, True, PAPER_DIM, MONO)])
    text(s, M, Inches(2.55), Inches(11), Inches(1.5),
         [("Ascend", 74, True, PAPER, SANS)])
    text(s, M, Inches(3.95), Inches(11), Inches(0.5),
         [("ON-PREM TO CLOUD DECISION SIMULATOR", 15, True, ACCENT, MONO)])
    text(s, M, Inches(4.75), Inches(9.2), Inches(1.1),
         [("A migration business case you can defend line by line — "
           "including the parts that argue against moving.", 17, False, PAPER_DIM,
           SANS)], spacing=1.35)
    rect(s, M, Inches(6.05), Inches(4.2), Inches(0.02), RGBColor(0x2A, 0x33, 0x4A))
    text(s, M, Inches(6.28), Inches(8), Inches(0.4),
         [(URL, 13, True, PAPER, MONO)])


def problem(prs, n):
    s = slide(prs)
    eyebrow(s, "The situation")
    heading(s, "The decision is being made anyway",
            "Renewal quotes have landed, budgets are being set, and the analysis that "
            "should inform the answer usually arrives after it.")
    y = Inches(2.95)
    cw = (CONTENT_W - Inches(0.6)) / 3
    card(s, M, y, cw, Inches(2.75), "Priced from a spreadsheet",
         "Cloud run-rate estimates built on list prices copied months ago, with no "
         "reserved-instance or licensing treatment.", NEGATIVE)
    card(s, M + cw + Inches(0.3), y, cw, Inches(2.75), "Sized from a guess",
         "Business cases assume a blanket 30-40% right-sizing saving. The estate rarely "
         "gives it back.", WARNING)
    card(s, M + 2 * (cw + Inches(0.3)), y, cw, Inches(2.75), "Scoped from a demo",
         "Tooling coverage taken at vendor face value, so the workloads it cannot carry "
         "surface during the programme, not before it.", ACCENT_2)
    footer(s, n)


def what_it_is(prs, n):
    s = slide(prs)
    eyebrow(s, "What Ascend is")
    heading(s, "One model, priced live, that answers the whole question",
            "Point it at an estate and it produces the run rate, the programme cost, the "
            "timeline, the risk and the business case — from a single set of assumptions "
            "you can see and change.")
    items = [
        ("Every Azure figure is priced against live Microsoft retail rates",
         "Fetched at run time from the Azure Retail Prices API for your region and "
         "currency, including reservations and savings plans. AWS comes from its own "
         "public rate feed."),
        (f"{FACTS['engines']} engines behind {FACTS['pages']} screens",
         "Right-sizing, readiness, 7R disposition, cost, complexity and effort, wave "
         "planning, Monte Carlo risk, TCO and NPV — one pipeline, so a change anywhere "
         "moves everything."),
        ("Organised as a seven-act client narrative, not a menu",
         "Situation, Discover, Assess, Decide, Plan, Execute, Communicate. Every screen "
         "states where it sits in the conversation."),
        ("Built to be presented",
         "A presentation mode hides every control, and each screen carries the one line "
         "worth saying out loud."),
    ]
    bullets(s, M, Inches(3.15), CONTENT_W, items, gap=Inches(0.92))
    footer(s, n)


def divider(prs, act, title, sub):
    s = slide(prs, dark=True)
    accent_bar(s)
    text(s, M, Inches(2.7), Inches(10), Inches(0.4),
         [(act.upper(), 13, True, ACCENT, MONO)])
    text(s, M, Inches(3.15), Inches(11), Inches(1.0),
         [(title, 44, True, PAPER, SANS)])
    text(s, M, Inches(4.35), Inches(9), Inches(0.8),
         [(sub, 15, False, PAPER_DIM, SANS)], spacing=1.3)


def findings(prs, n):
    s = slide(prs)
    eyebrow(s, "What the model says that a vendor will not")
    heading(s, "Four numbers that change the conversation",
            "All four are outputs of this model against a 547-VM reference estate. Two of "
            "them argue against the easy answer.")
    y, cw = Inches(2.9), (CONTENT_W - Inches(0.9)) / 4
    card(s, M, y, cw, Inches(3.25), "Azure and AWS charge the same",
         "Per vCPU hour for licence-included Windows, derived independently from both "
         "live APIs. The Azure case is Hybrid Benefit and free ESU — not cheaper cores.",
         ACCENT, big=FACTS["vcpu_rate"])
    card(s, M + cw + Inches(0.3), y, cw, Inches(3.25), "is all Azure Migrate carries",
         "End to end, and it moves no database of any kind. The rest needs tooling you "
         "have to choose, buy and staff.", WARNING, big=FACTS["migrate_coverage"])
    card(s, M + 2 * (cw + Inches(0.3)), y, cw, Inches(3.25),
         "is what right-sizing returns",
         f"Despite {FACTS['mean_cpu']} mean CPU. A 1.3x comfort factor and discrete SKU "
         "sizes eat the rest, so cases built on 30-40% do not survive the price list.",
         NEGATIVE, big=FACTS["rightsizing"])
    card(s, M + 3 * (cw + Inches(0.3)), y, cw, Inches(3.25),
         "hosts if you choose AVS",
         f"Sized by memory, not VM count, landing about {FACTS['avs_saving']} below "
         "staying on VMware. The fastest exit is rarely the cheapest home.",
         ACCENT_2, big=FACTS["avs_hosts"])
    footer(s, n)


def your_data(prs, n):
    s = slide(prs)
    eyebrow(s, "Using it on your estate")
    heading(s, "Your inventory in, in about two minutes",
            "Ascend refuses to show an answer until it knows whose estate it is modelling.")
    items = [
        ("Upload an RVTools export",
         "The vInfo sheet, or any CSV with a VM name, CPU, memory and guest OS. RVTools "
         "column names are recognised automatically."),
        ("It reads what the export actually carries",
         "Environment from your folder structure, workload tier and database engine from "
         "your naming conventions, applications from separated host names. Each one "
         "reports how many VMs it recognised."),
        ("It says plainly what it could not find",
         "vInfo has no performance counters, so right-sizing falls back to as-provisioned "
         "— the same position as a one-star Azure Migrate assessment. Ascend tells you, "
         "rather than presenting the gap as a finding."),
        ("No inventory yet? Enter six headline numbers",
         "VM count, Windows split, total vCPU, RAM, storage and current platform spend. "
         "The distribution is modelled; the totals are yours, and the app labels which "
         "half is which."),
    ]
    bullets(s, M, Inches(3.0), CONTENT_W, items, gap=Inches(0.95))
    footer(s, n)


def trust(prs, n):
    s = slide(prs)
    eyebrow(s, "Why the numbers can be defended")
    heading(s, "Measured, yours, or assumed — never blurred",
            "Every screen labels where its estate came from. A synthetic estate carries a "
            "warning on the executive briefing itself.")
    y = Inches(2.9)
    cw = (CONTENT_W - Inches(0.6)) / 3
    card(s, M, y, cw, Inches(2.9), "Live",
         "Azure Retail Prices API and the public AWS rate feed, for your region and "
         "currency, fetched while you watch. Google and Oracle use published list "
         "prices and are labelled as such.", POSITIVE)
    card(s, M + cw + Inches(0.3), y, cw, Inches(2.9), "Measured",
         "Your inventory: VM count, vCPU, RAM, storage, guest OS, clusters and hosts, "
         "exactly as discovery found them.", ACCENT)
    card(s, M + 2 * (cw + Inches(0.3)), y, cw, Inches(2.9), "Assumed, and listed",
         "Labour rates, replication bandwidth, discount rate and the on-premises cost "
         "base. Every one is on the Model & assumptions page with the screen that "
         "controls it.", WARNING)
    text(s, M, Inches(6.15), CONTENT_W, Inches(0.5),
         [("The two that drive every ratio are the estate and your current cost base. "
           "Supply both and the business case is yours, not ours.",
           13, True, INK, SANS)])
    footer(s, n)


def journey(prs, n):
    s = slide(prs)
    eyebrow(s, "The session")
    heading(s, "Seven acts, in the order a client asks the questions", "")
    acts = [
        ("Situation", "What is the answer, and what are we assuming?"),
        ("Discover", "What do you actually have?"),
        ("Assess", "What should happen to each workload, and how hard is it?"),
        ("Decide", "Which provider, which platforms, what does it cost, is it worth it?"),
        ("Plan", "How long, what constrains it, how confident should we be?"),
        ("Execute", "How does the tooling behave, and what else must you buy?"),
        ("Communicate", "Turning the model into deliverables."),
    ]
    y = Inches(2.5)
    for i, (act, q) in enumerate(acts):
        yy = y + i * Inches(0.62)
        rect(s, M, yy, Inches(0.42), Inches(0.42), ACCENT if i < 4 else ACCENT_2,
             MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, M, yy + Inches(0.08), Inches(0.42), Inches(0.3),
             [(str(i + 1), 13, True, PAPER, SANS)], align=PP_ALIGN.CENTER)
        text(s, M + Inches(0.62), yy + Inches(0.02), Inches(2.4), Inches(0.4),
             [(act, 15, True, INK, SANS)])
        text(s, M + Inches(3.1), yy + Inches(0.05), Inches(8.4), Inches(0.4),
             [(q, 13, False, INK_DIM, SANS)])
    footer(s, n)


def future_state(prs, n):
    s = slide(prs)
    eyebrow(s, "The picture the room asks for")
    heading(s, "Current state and future state, from the same model",
            "Not a drawing. Every line is a quantity the model computed, and it moves when "
            "an assumption moves.")
    y = Inches(2.95)
    cw = (CONTENT_W - Inches(0.5)) / 2
    card(s, M, y, cw, Inches(3.0), "Current state — VMware vSphere",
         "Clusters, hosts, datastores, virtual machines, Windows and Linux split, "
         "end-of-life guest operating systems, allocated vCPU and RAM, provisioned "
         "storage, and the host count the cost model derives from them.", NEGATIVE)
    card(s, M + cw + Inches(0.5), y, cw, Inches(3.0), "Future state — Azure",
         "Azure VMs and their SKU families, managed-service targets, Azure VMware "
         "Solution, workloads retired rather than migrated, right-sized vCPU, disk "
         "tiering, run rate, programme cost, waves and elapsed months.", POSITIVE)
    text(s, M, Inches(6.2), CONTENT_W, Inches(0.5),
         [("A flow diagram then carries every VM from its workload tier, through its 7R "
           "disposition, to the Azure service it lands on.", 13, True, INK, SANS)])
    footer(s, n)


def options(prs, n):
    """The fourteen destinations, grouped by the five kinds of answer they are."""
    s = slide(prs)
    eyebrow(s, "Leaving VMware")
    heading(s, "Fourteen destinations, six kinds of answer",
            "Ascend scores all fourteen, including the on-premises ones and doing nothing. "
            "The choice is not Azure or VMware; it is which of these six you are actually "
            "choosing between.")
    y = Inches(2.85)
    cw = (CONTENT_W - Inches(1.0)) / 3
    card(s, M, y, cw, Inches(1.95), "Azure public cloud",
         "Native IaaS, PaaS, and Azure VMware Solution. Leaves the data centre.", ACCENT)
    card(s, M + cw + Inches(0.5), y, cw, Inches(1.95), "Azure, in your data centre",
         "Azure Local, Azure Stack Hub, or Azure Arc managing the estate where it is.",
         ACCENT_2)
    card(s, M + 2 * (cw + Inches(0.5)), y, cw, Inches(1.95), "Another public cloud",
         "Google Cloud VMware Engine, Oracle Cloud VMware Solution and OCI native.",
         RGBColor(0x2E, 0x8F, 0xA3))
    y2 = y + Inches(2.2)
    card(s, M, y2, cw, Inches(1.95), "A different hypervisor, same racks",
         "Hyper-V with SCVMM, Nutanix AHV, OpenShift Virtualization, Proxmox VE. "
         "New licence, same building.", WARNING)
    card(s, M + cw + Inches(0.5), y2, cw, Inches(1.95), "Hybrid bare metal",
         "Nutanix Cloud Clusters on Azure: the same stack in a cloud region, on-premises, "
         "or both.", POSITIVE)
    card(s, M + 2 * (cw + Inches(0.5)), y2, cw, Inches(1.95),
         "Stay, and renegotiate",
         "VCF 9 on a renewed contract. Modelled as a real option, because sometimes it "
         "wins.", NEGATIVE)
    footer(s, n)


def which_wins(prs, n):
    """What the ranking engine says, and the fact that it changes with the goal."""
    s = slide(prs)
    eyebrow(s, "Which is better")
    heading(s, "It changes with what you are optimising for",
            "These are fit scores from Ascend's own ranking engine. Set the priority and "
            "the winner moves — which is the honest answer, and the reason the tool exists.")

    rows = [
        ("Exit the data centre on a deadline", "Azure native IaaS  80",
         "AVS 74 is the faster exit, but it is sized by memory and lands only ~5% below "
         "staying on VMware.", ACCENT),
        ("Cut total cost of ownership", "Azure native IaaS  80",
         "PaaS scores 76 and runs cheaper still, but costs far more to get to. "
         "Rehost first, replatform after.", POSITIVE),
        ("Escape Broadcom licensing", "Azure native IaaS  77",
         "Hyper-V with SCVMM scores 73 and keeps the workloads on your own racks, "
         "if leaving the building is not the point.", WARNING),
        ("Data must stay on-premises", "Hyper-V + SCVMM  83",
         "Azure Local scores 82 and gives an Azure control plane in your data centre. "
         "Public cloud is simply not an option here.", ACCENT_2),
    ]
    y = Inches(2.9)
    for i, (goal, winner, why, colour) in enumerate(rows):
        yy = y + i * Inches(0.92)
        rect(s, M, yy, Inches(0.05), Inches(0.62), colour)
        text(s, M + Inches(0.22), yy, Inches(3.5), Inches(0.4),
             [(goal, 13.5, True, INK, SANS)], spacing=1.15)
        text(s, M + Inches(3.95), yy, Inches(2.5), Inches(0.4),
             [(winner, 13.5, True, colour, SANS)])
        text(s, M + Inches(6.6), yy - Inches(0.02), Inches(4.9), Inches(0.7),
             [(why, 11.5, False, INK_DIM, SANS)], spacing=1.2)

    # Clear of the footer: the box must end above H - 0.55in or it strikes through it.
    rect(s, M, Inches(6.08), CONTENT_W, Inches(0.7), CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, M, Inches(6.08), Inches(0.06), Inches(0.7), ACCENT)
    text(s, M + Inches(0.3), Inches(6.2), CONTENT_W - Inches(0.6), Inches(0.5),
         [("Azure scores 92 against AWS 68, OCI 64 and Google 60 on a balanced weighting "
           "— but that lead is Hybrid Benefit and free ESU. Without Software Assurance it "
           "narrows sharply, and Ascend says so.", 12.5, True, INK, SANS)], spacing=1.15)
    footer(s, n)


def access(prs, n):
    s = slide(prs)
    eyebrow(s, "Access and data handling")
    heading(s, "Signed in, and nothing kept", "")
    items = [
        ("Accounts are issued per engagement",
         "PBKDF2-HMAC-SHA256 at 240,000 iterations with a per-account salt. Sessions "
         "expire after eight idle hours."),
        ("Your inventory stays in the session",
         "The uploaded file is held for the life of the browser session and written "
         "nowhere. Closing the tab discards it."),
        ("Nothing is sent to a vendor",
         "Pricing calls ask for public rate cards by region and SKU. No estate data "
         "leaves the application."),
        ("Everything is exportable",
         "Every table downloads as CSV, so the analysis can be reviewed in your own "
         "tooling rather than taken on trust."),
    ]
    bullets(s, M, Inches(2.5), CONTENT_W, items, gap=Inches(0.95))
    footer(s, n)


def next_steps(prs, n):
    s = slide(prs, dark=True)
    accent_bar(s)
    text(s, M, Inches(1.0), Inches(10), Inches(0.4),
         [("NEXT", 13, True, ACCENT, MONO)])
    text(s, M, Inches(1.45), Inches(11), Inches(0.9),
         [("What we need, and what you get back", 38, True, PAPER, SANS)])
    steps = [
        ("1", "An RVTools export", "The vInfo sheet from your vCenter. One file."),
        ("2", "Your current platform cost",
         "Total annual spend on running the VMware estate. A single number is enough; "
         "the split can be ours."),
        ("3", "A two-hour working session",
         "We walk the seven acts live, change assumptions in the room, and settle the "
         "ones that matter to you."),
        ("4", "A defensible business case",
         "Priced against live rates, with every assumption named and every table "
         "exportable."),
    ]
    y = Inches(2.95)
    for i, (num, head, body) in enumerate(steps):
        yy = y + i * Inches(0.95)
        rect(s, M, yy, Inches(0.52), Inches(0.52),
             ACCENT if i < 2 else ACCENT_2, MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, M, yy + Inches(0.11), Inches(0.52), Inches(0.35),
             [(num, 15, True, PAPER, SANS)], align=PP_ALIGN.CENTER)
        text(s, M + Inches(0.78), yy, Inches(3.5), Inches(0.4),
             [(head, 16, True, PAPER, SANS)])
        text(s, M + Inches(4.5), yy + Inches(0.03), Inches(7.4), Inches(0.7),
             [(body, 13, False, PAPER_DIM, SANS)], spacing=1.2)
    rect(s, M, Inches(6.75), Inches(4.2), Inches(0.02), RGBColor(0x2A, 0x33, 0x4A))
    text(s, M, Inches(6.95), Inches(8), Inches(0.4),
         [(URL, 13, True, PAPER, MONO)])


# --------------------------------------------------------------------------
def main() -> int:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    title_slide(prs)
    problem(prs, 2)
    what_it_is(prs, 3)
    divider(prs, "The evidence",
            "It will tell you when the answer is no",
            "A tool that only ever justifies the migration is a sales deck with "
            "arithmetic. This one is built to survive scrutiny.")
    findings(prs, 5)
    trust(prs, 6)
    divider(prs, "Using it",
            "On your estate, in your session",
            "Nothing is presented as an answer until Ascend knows whose estate it is "
            "modelling.")
    your_data(prs, 8)
    journey(prs, 9)
    future_state(prs, 10)
    options(prs, 11)
    which_wins(prs, 12)
    access(prs, 13)
    next_steps(prs, 14)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)

    print(f"Wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, "
          f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    for i, sl in enumerate(prs.slides, 1):
        titles = [sh.text_frame.paragraphs[0].runs[0].text
                  for sh in sl.shapes
                  if sh.has_text_frame and sh.text_frame.paragraphs[0].runs]
        print(f"  {i:2}. {titles[1] if len(titles) > 1 else titles[0] if titles else ''}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
